# text_extractor.py — UNIVERSAL EXTRACTOR (Enhanced)
# Improved handling for images, audio, video with better text extraction

import os, re, subprocess, tempfile
from pathlib import Path
from typing import Tuple, Dict, Any, List

from PIL import Image, ImageEnhance, ImageFilter, ImageOps
import pytesseract

if os.name == "nt":
    _tess = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    if os.path.exists(_tess):
        pytesseract.pytesseract.tesseract_cmd = _tess


def _lazy_imports():
    global fitz, PyPDF2, DocxDocument, openpyxl, Presentation, whisper, np
    try:
        import fitz  # PyMuPDF
    except Exception:
        fitz = None
    try:
        import PyPDF2
    except Exception:
        PyPDF2 = None
    try:
        from docx import Document as DocxDocument
    except Exception:
        DocxDocument = None
    try:
        import openpyxl
    except Exception:
        openpyxl = None
    try:
        from pptx import Presentation
    except Exception:
        Presentation = None
    try:
        import whisper  # openai-whisper
    except Exception:
        whisper = None
    try:
        import numpy as np
    except Exception:
        np = None


_lazy_imports()


# ---------- helpers ----------
def _enhance(img: Image.Image) -> Image.Image:
    """Enhanced image preprocessing for better OCR."""
    if img.mode != "L":
        img = img.convert("L")

    # Multiple enhancement strategies
    img = ImageFilter.MedianFilter(size=3)(img)
    img = ImageOps.autocontrast(img, cutoff=2)
    img = ImageEnhance.Contrast(img).enhance(2.0)
    img = ImageEnhance.Sharpness(img).enhance(1.5)

    # Additional enhancement for small text
    if img.size[0] < 1000 or img.size[1] < 1000:
        # Upscale small images
        new_size = (img.size[0] * 2, img.size[1] * 2)
        img = img.resize(new_size, Image.Resampling.LANCZOS)

    return img


def clean_text_spacing(text: str) -> str:
    text = text.replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def normalize_keywords(text: str) -> str:
    fixes = {
        r"obj\s*ective": "objective",
        r"intro\s*duction": "introduction",
        r"concl\s*usion": "conclusion",
        r"meth\s*odology": "methodology",
        r"depart\s*ment": "department",
        r"col\s*lege": "college",
        r"uni\s*versity": "university",
        r"mem\s*ber": "member",
        r"mem\s*bers": "members",
        r"tea\s*m": "team",
        r"na\s*me": "name",
    }
    for k, v in fixes.items():
        text = re.sub(k, v, text, flags=re.I)
    return text


def guess_type(path: str) -> str:
    ext = Path(path).suffix.lower()
    mapx = {
        ".pdf": "pdf", ".docx": "docx", ".doc": "docx",
        ".pptx": "pptx", ".ppt": "pptx",
        ".xlsx": "excel", ".xls": "excel", ".csv": "csv",
        ".txt": "text",
        ".jpg": "image", ".jpeg": "image", ".png": "image", ".tif": "image", ".tiff": "image", ".bmp": "image",
        ".webp": "image",
        ".mp3": "audio", ".wav": "audio", ".m4a": "audio", ".flac": "audio",
        ".mp4": "video", ".mov": "video", ".mkv": "video",
        ".py": "code", ".js": "code", ".json": "code", ".html": "code", ".xml": "code"
    }
    return mapx.get(ext, "unknown")


# ---------- Excel ----------
def _extract_excel(path: str) -> Tuple[str, Dict[str, Any]]:
    if openpyxl is None:
        return "[Excel support requires: pip install openpyxl]", {}
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        parts, meta = [], {"sheets": wb.sheetnames, "rows": 0}
        for name in wb.sheetnames:
            sh = wb[name]
            parts.append(f"\n=== Sheet: {name} ===\n")
            for row in sh.iter_rows(values_only=True):
                row_text = " | ".join("" if c is None else str(c) for c in row)
                if row_text.strip():
                    parts.append(row_text)
                    meta["rows"] += 1
        return "\n".join(parts), meta
    except Exception as e:
        return f"[Excel Error: {e}]", {}


# ---------- PPTX ----------
def _extract_pptx(path: str) -> Tuple[str, Dict[str, Any]]:
    if Presentation is None:
        return "[PowerPoint support requires: pip install python-pptx]", {}
    try:
        prs = Presentation(path)
        out, meta = [], {"slides": len(prs.slides), "has_notes": False}
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"\n=== Slide {i} ===\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    out.append(shape.text.strip())
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    out.append(f"[Notes] {notes}")
                    meta["has_notes"] = True
        return "\n".join(out), meta
    except Exception as e:
        return f"[PowerPoint Error: {e}]", {}


# ---------- PDF (text + layout + OCR fallback) ----------
def _extract_pdf(path: str) -> Tuple[str, Dict[str, Any]]:
    text = ""
    meta: Dict[str, Any] = {}

    # Try PyPDF2 text
    if PyPDF2 is not None:
        try:
            reader = PyPDF2.PdfReader(str(path))
            meta["pages"] = len(reader.pages)
            for p in reader.pages:
                text += (p.extract_text() or "") + "\n"
        except Exception:
            text = ""

    # If too short, try PyMuPDF layout blocks
    if (not text or len(text.strip()) < 80) and fitz is not None:
        try:
            doc = fitz.open(path)
            meta["pages"] = doc.page_count
            blocks_text: List[str] = []
            for i in range(doc.page_count):
                page = doc.load_page(i)
                for b in page.get_text("blocks"):
                    if len(b) >= 5:
                        blk = b[4].strip()
                        if blk:
                            blocks_text.append(blk)
            text = "\n\n".join(blocks_text)
        except Exception:
            pass

    # OCR fallback for scanned PDFs
    if (not text or len(text.strip()) < 80) and fitz is not None:
        try:
            doc = fitz.open(path)
            ocr_parts = []
            for i in range(doc.page_count):
                page = doc.load_page(i)
                pix = page.get_pixmap(dpi=200)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                img = _enhance(img)
                ocr_parts.append(pytesseract.image_to_string(img))
            text = "\n\n".join(ocr_parts)
            meta["ocr_applied"] = True
        except Exception as e:
            text = f"[PDF OCR Error: {e}]"

    return text, meta


# ---------- Image (Enhanced OCR) ----------
def _extract_image(path: str) -> Tuple[str, Dict[str, Any]]:
    """Enhanced image OCR with multiple passes and configuration."""
    try:
        img = Image.open(str(path))
        original_size = img.size

        # Try multiple OCR configurations for better results
        results = []

        # Config 1: Standard enhanced
        img_enhanced = _enhance(img.copy())
        text1 = pytesseract.image_to_string(img_enhanced, config='--psm 6')
        results.append(text1)

        # Config 2: Different page segmentation mode
        text2 = pytesseract.image_to_string(img_enhanced, config='--psm 3')
        results.append(text2)

        # Config 3: Original image (in case enhancement hurt)
        if img.mode != "RGB":
            img_rgb = img.convert("RGB")
        else:
            img_rgb = img
        text3 = pytesseract.image_to_string(img_rgb, config='--psm 6')
        results.append(text3)

        # Choose the longest result (usually most complete)
        text = max(results, key=len)

        # If still too short, add contextual description
        if len(text.strip()) < 20:
            text = f"[Image extracted from {Path(path).name}]\n\n{text}\n\n[Note: This appears to be an image with limited text content. Visual analysis may be needed.]"

        meta = {
            "image_size": original_size,
            "ocr_applied": True,
            "text_length": len(text)
        }

        return text, meta

    except Exception as e:
        return f"[Image OCR Error: {e}]", {"error": str(e)}


# ---------- Audio (Enhanced Transcription) ----------
def _extract_audio(path: str, whisper_model: str = "base") -> Tuple[str, Dict[str, Any]]:
    """Enhanced audio transcription with better formatting."""
    if whisper is None:
        return "[Audio Error: openai-whisper not installed]", {}
    try:
        model = whisper.load_model(whisper_model)
        res = model.transcribe(str(path), fp16=False, language="en")

        # Get the full text first
        full_text = res.get("text", "").strip()

        # Format with timestamps for reference
        segs = res.get("segments", [])
        timestamped_parts = []
        plain_parts = []

        for s in segs:
            t = s.get("text", "").strip()
            if t:
                start = s.get('start', 0)
                end = s.get('end', 0)
                timestamped_parts.append(f"[{start:.1f}s - {end:.1f}s] {t}")
                plain_parts.append(t)

        # Create comprehensive text that's searchable
        formatted_text = f"""[Audio Transcript]
File: {Path(path).name}
Duration: {res.get('duration', 0):.1f} seconds

Full Transcript:
{full_text}

Detailed Segments:
{chr(10).join(timestamped_parts)}"""

        return formatted_text, {
            "segments": len(segs),
            "duration": res.get("duration"),
            "language": res.get("language", "en")
        }
    except Exception as e:
        return f"[Audio Error: {e}]", {"error": str(e)}


# ---------- Video -> Audio (Enhanced) ----------
def _extract_video(path: str, whisper_model: str = "base") -> Tuple[str, Dict[str, Any]]:
    """Enhanced video transcription with better formatting."""
    if whisper is None:
        return "[Video Error: openai-whisper not installed]", {}
    try:
        tmp = tempfile.NamedTemporaryFile(suffix=".wav", delete=False).name
        cmd = ["ffmpeg", "-i", str(path), "-vn", "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1", tmp, "-y"]
        subprocess.run(cmd, capture_output=True, check=False)

        model = whisper.load_model(whisper_model)
        res = model.transcribe(tmp, fp16=False, language="en")

        try:
            os.unlink(tmp)
        except:
            pass

        # Get the full text
        full_text = res.get("text", "").strip()

        # Format segments
        segs = res.get("segments", [])
        timestamped_parts = []

        for s in segs:
            t = s.get("text", "").strip()
            if t:
                start = s.get('start', 0)
                end = s.get('end', 0)
                timestamped_parts.append(f"[{start:.1f}s - {end:.1f}s] {t}")

        # Create comprehensive text
        formatted_text = f"""[Video Audio Transcript]
File: {Path(path).name}
Duration: {res.get('duration', 0):.1f} seconds

Full Transcript:
{full_text}

Detailed Segments:
{chr(10).join(timestamped_parts)}"""

        return formatted_text, {
            "segments": len(segs),
            "duration": res.get("duration"),
            "language": res.get("language", "en")
        }
    except FileNotFoundError:
        return "[Video Error: ffmpeg not found in PATH]", {}
    except Exception as e:
        return f"[Video Error: {e}]", {"error": str(e)}


# ---------- Main ----------
def extract_text(path: str, whisper_model: str = "base") -> Tuple[str, Dict[str, Any]]:
    p = Path(path)
    typ = guess_type(str(p))
    meta: Dict[str, Any] = {"type": typ, "source": p.name, "size_bytes": p.stat().st_size}
    text = ""

    if typ == "pdf":
        text, extra = _extract_pdf(str(p))
        meta.update(extra)
    elif typ == "docx":
        if DocxDocument is None:
            text = "[DOCX support requires: pip install python-docx]"
        else:
            try:
                doc = DocxDocument(str(p))
                text = "\n".join(par.text for par in doc.paragraphs if par.text.strip())
                for table in getattr(doc, "tables", []):
                    for row in table.rows:
                        row_text = " | ".join(c.text.strip() for c in row.cells)
                        if row_text.strip():
                            text += "\n" + row_text
            except Exception as e:
                text = f"[DOCX Error: {e}]"
    elif typ == "pptx":
        text, extra = _extract_pptx(str(p))
        meta.update(extra)
    elif typ == "excel":
        text, extra = _extract_excel(str(p))
        meta.update(extra)
    elif typ == "image":
        text, extra = _extract_image(str(p))
        meta.update(extra)
    elif typ == "audio":
        text, extra = _extract_audio(str(p), whisper_model=whisper_model)
        meta.update(extra)
    elif typ == "video":
        text, extra = _extract_video(str(p), whisper_model=whisper_model)
        meta.update(extra)
    elif typ in ["text", "csv", "code"]:
        tried = ["utf-8", "utf-16", "latin-1", "cp1252"]
        for enc in tried:
            try:
                with open(p, "r", encoding=enc, errors="ignore") as f:
                    text = f.read()
                meta["encoding"] = enc
                break
            except Exception:
                continue
        if not text:
            text = "[Text Read Error]"
    else:
        text = "[Unsupported file type]"

    if text and not text.startswith("["):
        text = clean_text_spacing(text)
        text = normalize_keywords(text)
        meta["final_char_count"] = len(text)

    return text, meta